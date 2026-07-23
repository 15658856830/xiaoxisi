import os, sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

BG     = RGBColor(0xF7, 0xFB, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEAL   = RGBColor(0x00, 0x8F, 0x7B)
BLUE   = RGBColor(0x00, 0x9E, 0xDB)
GREEN  = RGBColor(0x72, 0xB5, 0x48)
ORANGE = RGBColor(0xD9, 0x85, 0x00)
CC     = [TEAL, BLUE, GREEN, ORANGE]
SW     = 12198350  # slide width in EMU (标准 16:9 幻灯片)
SH     = 6859270   # slide height in EMU
BH     = 857449    # banner height in EMU (顶部色条)


def rect(slide, x, y, w, h, fill):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    r.line.fill.background()
    return r._element


def main(inp, out):
    if not os.path.exists(inp):
        print(f"输入文件不存在: {inp}")
        sys.exit(1)
    prs = Presentation(inp)
    print(f"reading {len(prs.slides)} slides...")
    
    for i, slide in enumerate(prs.slides, 1):
        tree = slide.shapes._spTree
        elms = []
        
        # 1) Background (z=0)
        elms.append(rect(slide, 0, 0, SW, SH, BG))
        # 2) Banner (z=1)
        elms.append(rect(slide, 0, 0, SW, BH, TEAL))
        
        if 2 <= i <= 18:
            # 3) Section title bar (z=2)
            for s in slide.shapes:
                if not s.has_text_frame:
                    continue
                t = s.text_frame.text.strip()
                if t[:2] in ("一、", "二、"):
                    elms.append(rect(slide, 0, s.top-30000, SW, s.height+60000, TEAL))
                    break
        
        # 4) Text cards for text-heavy slides
        imgs = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        charts = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.CHART]
        
        if len(imgs) <= 3 and len(charts) == 0 and i not in (1, 19) and i != 2:
            txts = []
            for s in slide.shapes:
                if not s.has_text_frame:
                    continue
                t = s.text_frame.text.strip()
                if not t or t[:2] in ("\u4e00\u3001", "\u4e8c\u3001"):
                    continue
                txts.append(s)
            
            if len(txts) >= 2:
                for ci, ts in enumerate(txts):
                    pad = 35000
                    cx = max(ts.left - pad, 60000)
                    cy = ts.top - pad
                    cw = min(ts.width + pad*2, SW - 200000)
                    ch = ts.height + pad*2
                    elms.append(rect(slide, cx, cy, cw, ch, WHITE))
                    elms.append(rect(slide, cx, cy, 48000, ch, CC[ci % len(CC)]))
        
        # Move all new elements to back (z=0,1,2,3...)
        for el in elms:
            tree.remove(el)
        for ei, el in enumerate(elms):
            tree.insert(ei, el)
        
        n_new = len(elms)
        print(f"  P{i:2d} +{n_new} new shapes (total now: {len(list(slide.shapes))})")
    
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"\ndone: {out}")


if __name__ == '__main__':
    main(sys.argv[1], sys.argv[2])
