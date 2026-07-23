"""Polish presentation text - fix grammar, punctuation, flow. Keep data intact."""
import os, sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

CN_FONT = "微软雅黑"
EN_FONT = "Arial"
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)


def has_cn(text):
    return any('\u4e00' <= c <= '\u9fff' for c in text)


def pick_font(text):
    return CN_FONT if has_cn(text) else EN_FONT


def replace_para_text(para, new_text):
    """Replace paragraph text while maintaining font."""
    para.clear()
    run = para.add_run()
    run.text = new_text
    run.font.name = pick_font(new_text)
    run.font.size = Pt(10)
    run.font.color.rgb = DARK_TEXT


# replacement rules: (slide_index, original_substring, new_text)
# slide_index 0 means match any slide
REPLACEMENTS = [
    # P8 - language polish
    (8,  "今年总共培养11名研究生，其中博士8名，硕士3名",
          "本年度共培养研究生11名（博士8名、硕士3名）"),
    # P9 - language polish
    (9,  "逐步构建成上海区域内有影响力的开放性、友好型的血管病科研转化平台",
          "逐步构建上海区域内有影响力的开放性、友好型血管病科研转化平台"),
    (9,  "目前已经完成11万的横向课题入院。",
          "目前已完成11万元横向课题入院。"),
    (9,  "2026年计划完成20万+横向课题入院。",
          "2026年计划完成20万元以上横向课题入院。"),
    # P10 - trailing comma
    (10, "目前在研课题共9项、",
          "目前在研课题共9项"),
    (10, "今年发表文章IF共计5.5",
          "本年度发表文章IF共计5.5"),
    # P11 - fix quotes, split into natural items
    (11, "荣获上海市科学技术普及奖二等奖、2026年度医院\u201c先进共产党员、在2025年度工作中做出突出贡献",
          "荣获上海市科学技术普及奖二等奖\n2026年度医院\u201c先进共产党员\u201d\n在2025年度工作中做出突出贡献"),
    # P16 - minor flow polish
    (16, "受限于浦东医院等级，与传统强院相比，在进修生来源、对外影响力及受众信任度有差异；",
          "受限于浦东医院等级，与传统强院相比，在进修生来源、对外影响力及受众信任度方面存在差异；"),
    (16, "因血管外科手术费收费的低廉、耗材费的高昂，将会受到DRGS/DIP政策的严重影响。",
          "因血管外科手术收费标准低、耗材费用高，将受到DRGS/DIP政策的较大影响。"),
]


def main(in_path, out_path):
    if not os.path.exists(in_path):
        print(f"输入文件不存在: {in_path}")
        sys.exit(1)
    print("reading...")
    prs = Presentation(in_path)
    changes = 0

    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                full = para.text.strip()
                if not full:
                    continue
                for sidx, orig, new in REPLACEMENTS:
                    if sidx != i and sidx != 0:
                        continue
                    if orig in full:
                        new_text = full.replace(orig, new)
                        replace_para_text(para, new_text)
                        changes += 1
                        print(f"  P{i}: replaced \"{orig[:40]}...\"")

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print(f"\ndone: {out_path}")
    print(f"  total changes: {changes}")


if __name__ == '__main__':
    main(sys.argv[1], sys.argv[2])
