"""PPT beautify - fonts, tables, charts"""
import os, sys
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# color scheme
PRIMARY    = RGBColor(0x2B, 0x57, 0x9A)
SECONDARY  = RGBColor(0x5B, 0x9B, 0xD5)
ACCENT     = RGBColor(0x70, 0xAD, 0x47)
WARM       = RGBColor(0xED, 0x7D, 0x31)
LIGHT_BG   = RGBColor(0xD6, 0xE4, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)

CHART_COLORS = [
    RGBColor(0x2B, 0x57, 0x9A),
    RGBColor(0x5B, 0x9B, 0xD5),
    RGBColor(0x70, 0xAD, 0x47),
    RGBColor(0xED, 0x7D, 0x31),
    RGBColor(0xA5, 0xA5, 0xA5),
    RGBColor(0xFF, 0xC0, 0x00),
]

CN_FONT = "微软雅黑"
EN_FONT = "Arial"


def contains_cn(text):
    return any('\u4e00' <= c <= '\u9fff' for c in text)


def set_cell_fill(cell, color):
    tcPr = cell._tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))


def set_cell_margin(cell, pt=2):
    tcPr = cell._tc.get_or_add_tcPr()
    emu = str(int(Pt(pt)))
    for a in ('marL', 'marR', 'marT', 'marB'):
        tcPr.set(a, emu)


def set_run_font(run, sz=None, bold=False, color=None, fn=None):
    f = run.font
    if sz:
        f.size = Pt(sz)
    f.bold = bold
    if color:
        f.color.rgb = color
    if fn:
        f.name = fn


def beautify_fonts(slide):
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        for p in s.text_frame.paragraphs:
            for r in p.runs:
                if not r.text.strip():
                    continue
                r.font.name = CN_FONT if contains_cn(r.text) else EN_FONT


def beautify_table(table):
    nrows = len(table.rows)
    ncols = len(table.columns)
    for ri in range(nrows):
        for ci in range(ncols):
            cell = table.cell(ri, ci)
            if ri == 0:
                set_cell_fill(cell, PRIMARY)
                for p in cell.text_frame.paragraphs:
                    p.alignment = 1
                    for r in p.runs:
                        set_run_font(r, sz=10, bold=True, color=WHITE)
            else:
                set_cell_fill(cell, LIGHT_BG if ri % 2 == 0 else WHITE)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        set_run_font(r, sz=9, color=DARK_TEXT)
            set_cell_margin(cell, 2)


def beautify_chart(chart):
    try:
        chart.chart_style = 2
    except Exception:
        pass
    for plot in chart.plots:
        for si, series in enumerate(plot.series):
            c = CHART_COLORS[si % len(CHART_COLORS)]
            try:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = c
            except Exception:
                pass
        if hasattr(plot, 'has_data_labels') and plot.has_data_labels:
            try:
                dl = plot.data_labels
                dl.font.size = Pt(9)
                dl.font.color.rgb = DARK_TEXT
            except Exception:
                pass


def main(in_path, out_path):
    if not os.path.exists(in_path):
        print(f"输入文件不存在: {in_path}")
        sys.exit(1)
    print("reading...")
    prs = Presentation(in_path)
    st = {'pages': 0, 'tables': 0, 'charts': 0}
    for i, slide in enumerate(prs.slides, 1):
        beautify_fonts(slide)
        st['pages'] += 1
        for s in slide.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.TABLE:
                beautify_table(s.table)
                st['tables'] += 1
            if s.shape_type == MSO_SHAPE_TYPE.CHART:
                beautify_chart(s.chart)
                st['charts'] += 1
        print(f"  slide {i:2d} done")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print(f"done: {out_path}")
    print(f"  pages: {st['pages']}, tables: {st['tables']}, charts: {st['charts']}")


if __name__ == '__main__':
    main(sys.argv[1], sys.argv[2])
